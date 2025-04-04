import os
import tempfile
import unittest
import datetime

from pptx import Presentation
from pptx.util import Inches

from template_reports.office_renderer import render_pptx


# Dummy objects for integration testing.
class DummyUser:
    def __init__(self, name, email, is_active=True):
        self.name = name
        self.email = email
        self.is_active = is_active
        self._meta = True  # Simulate a Django model

    def __str__(self):
        return self.name


class DummyCohort:
    def __init__(self, name):
        self.name = name

    def __str__(self):
        return self.name


class DummyRequestUser:
    def has_perm(self, perm, obj):
        # Deny permission if object's name contains "deny"
        if hasattr(obj, "name") and "deny" in obj.name.lower():
            return False
        return True


class TestRendererIntegration(unittest.TestCase):
    def setUp(self):
        # Create a minimal PPTX with one text box and one table.
        self.prs = Presentation()
        blank_slide = self.prs.slide_layouts[5]
        self.slide = self.prs.slides.add_slide(blank_slide)
        # Add a text box with mixed text.
        textbox = self.slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4), Inches(1)
        )
        textbox.text_frame.text = "Welcome, {{ user.name }}. Program: {{ program.name }}."
        self.textbox_index = len(self.slide.shapes) - 1
        # Add a table with one row, one column.
        rows, cols = 1, 1
        left, top, width, height = Inches(0.5), Inches(2), Inches(4), Inches(0.8)
        table_shape = self.slide.shapes.add_table(rows, cols, left, top, width, height)
        # Use the original mixed text with prefix.
        table_cell = table_shape.table.cell(0, 0)
        table_cell.text = "Here: {{ program.users.email }}"
        self.table_shape_index = None
        for idx, shape in enumerate(self.slide.shapes):
            if getattr(shape, "has_table", False) and shape.has_table:
                self.table_shape_index = idx
                break

        # Save this PPTX to a temporary file.
        self.temp_input = tempfile.mktemp(suffix=".pptx")
        self.temp_output = tempfile.mktemp(suffix=".pptx")
        self.prs.save(self.temp_input)

        # Set up the context.
        self.cohort = DummyCohort("Cohort A")
        self.user = DummyUser("Alice", "alice@example.com", is_active=True)
        self.user2 = DummyUser("Bob", "bob@example.com", is_active=True)
        self.user3 = DummyUser("Carol", "carol@example.com", is_active=False)
        self.program = {
            "name": "Test Program",
            "users": [self.user, self.user2, self.user3],
        }
        self.context = {
            "user": self.user,
            "program": self.program,
            "date": datetime.date(2020, 1, 15),
        }
        self.request_user = DummyRequestUser()

    def tearDown(self):
        if os.path.exists(self.temp_input):
            os.remove(self.temp_input)
        if os.path.exists(self.temp_output):
            os.remove(self.temp_output)

    def test_integration_renderer(self):
        # Run the renderer integration.
        rendered, errors = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            perm_user=None,
        )
        self.assertIsNone(errors)

        # Open the rendered PPTX.
        prs_out = Presentation(rendered)
        # Test text box content.
        textbox = prs_out.slides[0].shapes[self.textbox_index]
        txt = textbox.text_frame.text
        self.assertIn("Welcome, Alice.", txt)
        self.assertIn("Program: Test Program", txt)
        # Check table content.
        table_shape = prs_out.slides[0].shapes[self.table_shape_index]
        self.assertIsNotNone(table_shape)
        # The original cell text is "Here: {{ program.users.email }}", so after rendering
        # it should expand into as many rows as there are program.users.
        expected_users = [u.email for u in self.program["users"]]
        # Get all rows from the table.
        rows = list(table_shape.table.rows)
        self.assertEqual(len(rows), len(expected_users))
        # For each row, check that its cell text equals "Here: " and then the corresponding email.
        for i, row in enumerate(rows):
            cell_text = row.cells[0].text.strip()
            expected = f"Here: {expected_users[i]}"
            self.assertEqual(cell_text, expected)


if __name__ == "__main__":
    unittest.main()
