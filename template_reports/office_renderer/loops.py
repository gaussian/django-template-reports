"""
Functions to handle loop functionality in PPTX templates.
"""

from __future__ import annotations

import re
from typing import TYPE_CHECKING, Iterable, Optional, Any

from ..templating import resolve_tag
from .constants import (
    LOOP_START_PATTERN_STR,
    LOOP_END_PATTERN_STR,
)
from .pptx_utils import duplicate_slide, remove_shape

if TYPE_CHECKING:
    from pptx.presentation import Presentation

# Patterns for loop directives
LOOP_START_PATTERN = re.compile(LOOP_START_PATTERN_STR)
LOOP_END_PATTERN = re.compile(LOOP_END_PATTERN_STR)


def extract_loop_directive(text: str | None) -> tuple[str | None, str | None]:
    """Return (variable, collection) if *text* contains a loop directive."""
    if not text:
        return None, None

    match = LOOP_START_PATTERN.search(text.strip())
    if match:
        variable = match.group(1)
        collection = match.group(2)
        return variable, collection

    return None, None


def is_loop_start(shape) -> bool:
    """Return True if the shape text indicates a loop start."""
    if not hasattr(shape, "text_frame") or not hasattr(shape.text_frame, "text"):
        return False

    variable, collection = extract_loop_directive(shape.text_frame.text)
    return variable is not None and collection is not None


def is_loop_end(shape) -> bool:
    """Return True if the shape text indicates a loop end."""
    if not hasattr(shape, "text_frame") or not hasattr(shape.text_frame, "text"):
        return False

    text = shape.text_frame.text
    if not text:
        return False

    return LOOP_END_PATTERN.search(text.strip()) is not None


def get_collection_from_collection_tag(
    collection_tag: str,
    context: dict,
    perm_user,
) -> tuple[Optional[Iterable[Any]], Optional[str]]:
    """
    Given a collection tag, resolve it to a collection in the context.
    This function is a placeholder and should be implemented based on your context resolution logic.
    """
    # Get the collection from the context using resolve_tag
    try:
        collection = resolve_tag(collection_tag, context, perm_user)
    except Exception as e:
        return (
            None,
            f"Failed to resolve collection '{collection_tag}': {str(e)}",
        )

    # No collection
    if collection is None:
        return None, f"Collection '{collection_tag}' not found in context"

    # Ensure the collection is iterable (force evaluation of any lazy iterables)
    try:
        collection = list(collection)
    except TypeError:
        return None, f"'{collection_tag}' is not a collection of things"

    # For empty or nonexistent collections
    if not collection:
        return None, f"Collection '{collection_tag}' is empty"

    return collection, None


def process_loops(prs: Presentation, context, perm_user, errors):
    """
    Process loops in the presentation:
    - Identify loop sections (slides between %loop var in collection% and %endloop%)
    - For each item in the collection, create a mapping between loop slides and variables
    - Return a list of slides to process with their context info
    """

    # First pass: identify loop sections and collect info about them
    slide_sections = []
    in_loop = False
    loop_variable_name = None
    loop_collection = None
    loop_slides = []

    for i, slide in enumerate(prs.slides):
        # Track whether this slide has loop directives
        has_loop_start = False
        has_loop_end = False
        loop_start_shapes = []
        loop_end_shapes = []

        for shape in slide.shapes:
            # Check for loop start
            if is_loop_start(shape):
                # Store the shape, but delete it from the slide
                loop_start_shapes.append(shape)
                remove_shape(shape)

                # Quit if multiple loop start directives on the same slide
                if len(loop_start_shapes) > 1:
                    errors.append(
                        f"Error on slide {i + 1}: Multiple loop start directives on same slide"
                    )
                    return []

                # Get the loop variable and collection data
                loop_variable_name, collection_tag = extract_loop_directive(
                    shape.text_frame.text
                )
                if loop_variable_name and collection_tag:
                    has_loop_start = True
                    loop_collection, error = get_collection_from_collection_tag(
                        collection_tag, context, perm_user
                    )
                    if error:
                        errors.append(f"Error on slide {i + 1}: {error}")
                        return []
                else:
                    errors.append(f"Error on slide {i + 1}: Invalid loop start directive")
                    return []

            # Check for loop end
            if is_loop_end(shape):
                # Store the shape, but delete it from the slide
                loop_end_shapes.append(shape)
                remove_shape(shape)

                # Quit if multiple loop end directives on the same slide
                if len(loop_end_shapes) > 1:
                    errors.append(
                        f"Error on slide {i + 1}: Multiple loop end directives on same slide"
                    )
                    return []

                if not in_loop and not has_loop_start:
                    errors.append(
                        f"Error on slide {i + 1}: %endloop% without a matching loop start"
                    )
                has_loop_end = True

        # Handle loop start
        if has_loop_start:
            if in_loop:
                errors.append(f"Error on slide {i + 1}: Nested loops are not supported")
                return []

            in_loop = True

        # Handle if in loop
        if in_loop:

            # Store slides, either in loop or not
            loop_slides.append(slide)

            # Handle loop end
            if has_loop_end:
                # Quit if loop end directives there are on the slide
                loop_end_count = sum(1 for shape in slide.shapes if is_loop_end(shape))
                if loop_end_count > 1:
                    errors.append(
                        f"Error on slide {i + 1}: Multiple loop end directives on same slide"
                    )
                    return []

                slide_sections.append(
                    {
                        "slides": loop_slides,
                        "loop_variable_name": loop_variable_name,
                        "loop_collection": loop_collection,
                    }
                )

                # Reset loop state
                in_loop = False
                loop_variable_name = None
                loop_collection = None
                loop_slides = []

        # Non-loop slides
        else:
            slide_sections.append(
                {
                    "slides": [slide],
                }
            )

        # # Handle loop slides
        # if in_loop:
        #     # For each item in the collection, process each slide
        #     assert collection
        #     for loop_item in collection:
        #         # Add the loop variable to the context for this slide
        #         extra_context = {loop_variable_name: loop_item}
        #         for j in range(section["start_index"], section["end_index"] + 1):
        #             # Duplicate the slide scenarios
        #             new_slide = duplicate_slide(prs, j)
        #             # Store the new slide in the list
        #             slides_to_process.append(
        #                 {
        #                     "slide": new_slide,
        #                     "slide_number": current_slide_number,
        #                     "extra_context": extra_context,
        #                 }
        #             )
        #             current_slide_number += 1

        # # Handle non-loop slides
        # else:
        #     slides_to_process.append(
        #         {
        #             "slide": slide,
        #             "slide_number": current_slide_number,
        #         }
        #     )
        #     current_slide_number += 1

    # Check for unclosed loops
    if in_loop:
        errors.append(f"Error: Loop started but never closed with %endloop%")
        return []  # Short-circuit when unclosed loop found

    # Keep track of which slides are part of loops to avoid duplicating them
    # loop_slide_indices = set()
    # for section in slide_sections:
    #     for i in range(section["start_index"], section["end_index"] + 1):
    #         loop_slide_indices.add(i)

    # Prepare slides to process
    slides_to_process = []
    current_slide_number = 1

    # Create the slide structure (incl duplication and extra loop context)
    for section in slide_sections:
        loop_variable_name = section.get("loop_variable_name", None)
        loop_collection = section.get("loop_collection", None)
        slides = section.get("slides")

        # Not in loop
        if not loop_variable_name:
            for slide in slides:
                slides_to_process.append(
                    {
                        "slide": slide,
                        "slide_number": current_slide_number,
                    }
                )
                current_slide_number += 1

        # In loop
        else:
            assert loop_collection is not None

            # For each item in the collection, process each slide
            loop_count = len(loop_collection)
            for i, loop_item in enumerate(loop_collection):
                # Add the loop variable to the context for this slide
                extra_context = {
                    section["loop_variable_name"]: loop_item,
                    "loop_count": loop_count,
                }
                for slide in slides:
                    # Duplicate the slide, or (if this is first loop item) use the original
                    if i == 0:
                        new_slide = slide
                    else:
                        # Current slide number is 1-indexed
                        new_slide_index = current_slide_number - 1
                        new_slide = duplicate_slide(prs, slide, new_slide_index)
                    # Store the new slide in the list
                    slides_to_process.append(
                        {
                            "slide": new_slide,
                            "slide_number": current_slide_number,
                            "extra_context": {
                                **extra_context,
                                "loop_number": i + 1,
                            },
                        }
                    )
                    current_slide_number += 1

    return slides_to_process
