import re
from pptx import Presentation
from .paragraphs import merge_split_placeholders

# Pattern to match placeholders, e.g. "{{ some.placeholder }}"
PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*(.*?)\s*\}\}")


def extract_top_level_context_keys_from_text(text):
    """
    Given a text string, find all placeholders and extract the top-level context key.
    E.g., from "users.programs.email" or "programs[is_active=True].name" return "users" and "programs".
    """
    keys = set()
    placeholders = PLACEHOLDER_PATTERN.findall(text)
    for ph in placeholders:
        ph = ph.strip()
        if ph:
            # Get first token before a dot or '['.
            m = re.match(r"([^\.\[\]]+)", ph)
            if m:
                keys.add(m.group(1))
    return keys


def extract_context_keys(template_path):
    """
    Iterate through all slides, shapes, paragraphs and table cells in the PPTX at template_path,
    merging split placeholders, and return a sorted list of unique top-level context keys found.
    """
    prs = Presentation(template_path)
    keys = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            # Process text frames.
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    merge_split_placeholders(paragraph)
                    keys.update(extract_top_level_context_keys_from_text(paragraph.text))
            # Process table cells.
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                merge_split_placeholders(paragraph)
                                keys.update(
                                    extract_top_level_context_keys_from_text(
                                        paragraph.text
                                    )
                                )
    return sorted(keys)
