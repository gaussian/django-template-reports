from pptx import Presentation
from .paragraphs import process_paragraph
from .tables import process_table_cell


def render_pptx(
    template_path, context, output_path, request_user=None, check_permissions=True
):
    """
    Render the PPTX template at template_path using the provided context and save to output_path.

    1) For non-table shapes (text frames):
       - We iterate over each paragraph and merge runs if a single placeholder is split across them,
         calling process_paragraph(...) in "normal" mode.
       - After merging, each run's text is replaced by process_text(..., mode="normal").

    2) For table shapes:
       - We iterate over each cell. If the entire cell is exactly a single placeholder, we call
         process_table_cell(...) in "table" mode. That function handles pure placeholders vs. mixed text.
       - If not a pure placeholder, we perform normal merges on each paragraph (so that placeholders
         split across runs are unified) and then call "normal" mode on the resulting text.

    After processing, if any errors were recorded:
      - If any contain "Permission denied", raise a PermissionDeniedException.
      - Otherwise, raise an UnresolvedTagError.

    Args:
      template_path (str): Path to the input PPTX template.
      context (dict): The template context.
      output_path (str): Path to save the rendered PPTX.
      request_user: The user object for permission checking.
      check_permissions (bool): Whether to enforce permission checks.

    Returns:
      The output_path if rendering succeeds.

    Raises:
      PermissionDeniedException, UnresolvedTagError, UnterminatedTagException.
    """
    prs = Presentation(template_path)
    errors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # 1) Process text frames (non-table).
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    # Merge any placeholders that are split across multiple runs.
                    process_paragraph(
                        paragraph,
                        context,
                        request_user,
                        check_permissions,
                        mode="normal",  # for text frames
                    )
            # 2) Process tables.
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process_table_cell(
                            cell, context, errors, request_user, check_permissions
                        )

    if errors:
        print("Rendering aborted due to the following errors:")
        for err in set(errors):
            print(f" - {err}")
        print("Output file not saved.")
        return None

    prs.save(output_path)
    return output_path
