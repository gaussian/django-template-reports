import os
from pptx import Presentation
from .templating import process_text


class UnresolvedTagError(Exception):
    """Raised when one or more template tags could not be resolved."""

    pass


class UnterminatedTagException(Exception):
    """Raised when a template tag starting with '{{' is not terminated by '}}' in the same paragraph."""

    pass


def merge_runs_in_paragraph(paragraph, context, errors, request_user, check_permissions):
    """
    Scans the runs in the paragraph and, if a run contains a starting '{{' without a closing '}}',
    merges subsequent runs into the first run until a closing '}}' is found.

    If no closing '}}' is found in the paragraph, raises UnterminatedTagException.

    After merging, the text is processed via process_text.

    This simply concatenates the text from the subsequent runs into the first run and
    deletes those subsequent runs.
    """
    runs = paragraph.runs
    i = 0
    while i < len(runs):
        current_text = runs[i].text
        if "{{" in current_text and "}}" not in current_text:
            merged_text = current_text
            j = i + 1
            found_closing = False
            while j < len(runs):
                merged_text += runs[j].text
                if "}}" in runs[j].text:
                    found_closing = True
                    break
                j += 1
            if not found_closing:
                raise UnterminatedTagException(
                    f"Unterminated tag starting in run {i}: {current_text}"
                )
            processed_text = process_text(
                merged_text,
                context=context,
                errors=errors,
                request_user=request_user,
                check_permissions=check_permissions,
            )
            # Replace text in the first run with processed text.
            runs[i].text = processed_text
            # Remove the merged runs from the paragraph using the underlying XML elements.
            for k in range(i + 1, j + 1):
                paragraph._p.remove(runs[k]._r)
            # Update the runs list by re-fetching paragraph.runs.
            # (This ensures we have an updated list after removal.)
            runs = paragraph.runs
        else:
            processed_text = process_text(
                current_text,
                context=context,
                errors=errors,
                request_user=request_user,
                check_permissions=check_permissions,
            )
            runs[i].text = processed_text
        i += 1


def render_pptx(
    template_path, context, output_path, request_user=None, check_permissions=True
):
    """
    Open the PPTX template at `template_path`, process all text placeholders (merging runs as needed)
    using our templating logic, and save the rendered presentation to `output_path`.

    If any tags cannot be resolved, prints them and raises UnresolvedTagError.
    If a tag is unterminated (no closing '}}' in a paragraph), raises UnterminatedTagException.
    """
    prs = Presentation(template_path)
    errors = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                merge_runs_in_paragraph(
                    paragraph, context, errors, request_user, check_permissions
                )
    if errors:
        print(
            "The following template tags could not be resolved or failed permission checks:"
        )
        for tag in set(errors):
            print(f" - {tag}")
        raise UnresolvedTagError(
            "One or more template tags could not be resolved; output file not created."
        )
    prs.save(output_path)
    return output_path
